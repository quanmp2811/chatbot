"""Utilities for extracting normalized UTF-8 text from local and Google Drive files."""

from __future__ import annotations

import io
import json
import re
import unicodedata
from pathlib import Path
from typing import Any, Callable
from xml.etree import ElementTree

import pandas as pd
from googleapiclient.http import MediaIoBaseDownload

try:
    import win32com.client
    word_app = None
except ImportError:
    win32com = None

try:
    from ebooklib import epub
except ImportError:  # pragma: no cover
    epub = None

try:
    import docx2txt
except ImportError:  # pragma: no cover
    docx2txt = None

try:
    from odf import teletype
    from odf.opendocument import load as load_odt
    from odf.text import H, P
except ImportError:  # pragma: no cover
    teletype = None
    load_odt = None
    H = None
    P = None


SUPPORTED_EXTENSIONS: dict[str, Callable[[str | Path], str]] = {
    ".pdf": lambda file_path: read_pdf(file_path),
    ".doc": lambda file_path: read_doc(file_path),
    ".docx": lambda file_path: read_docx(file_path),
    ".txt": lambda file_path: read_txt(file_path),
    ".csv": lambda file_path: read_csv(file_path),
    ".xlsx": lambda file_path: read_excel(file_path),
    ".pptx": lambda file_path: read_pptx(file_path),
    ".html": lambda file_path: read_html(file_path),
    ".htm": lambda file_path: read_html(file_path),
    ".json": lambda file_path: read_json(file_path),
    ".md": lambda file_path: read_markdown(file_path),
    ".markdown": lambda file_path: read_markdown(file_path),
    ".xml": lambda file_path: read_xml(file_path),
    ".log": lambda file_path: read_log(file_path),
    ".rtf": lambda file_path: read_rtf(file_path),
    ".odt": lambda file_path: read_odt(file_path),
    ".epub": lambda file_path: read_epub(file_path),
    ".png": lambda file_path: read_image_ocr(file_path),
    ".jpg": lambda file_path: read_image_ocr(file_path),
    ".jpeg": lambda file_path: read_image_ocr(file_path),
}

GOOGLE_DOC_MIME = "application/vnd.google-apps.document"
GOOGLE_SHEET_MIME = "application/vnd.google-apps.spreadsheet"
GOOGLE_SLIDE_MIME = "application/vnd.google-apps.presentation"


def _ensure_path(file_path: str | Path) -> Path:
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")
    if not path.is_file():
        raise ValueError(f"Path is not a file: {path}")
    return path


def _read_text_file(file_path: str | Path) -> str:
    path = _ensure_path(file_path)
    return path.read_text(encoding="utf-8", errors="ignore")


def _normalize_text(text: str) -> str:
    if not text:
        return ""

    text = unicodedata.normalize("NFKC", text)
    text = text.replace("\ufeff", "")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[^\S\n]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = "\n".join(line.strip() for line in text.split("\n"))
    text = re.sub(r" +", " ", text)
    return text.encode("utf-8", errors="ignore").decode("utf-8").strip()


def _dataframe_to_text(dataframe: pd.DataFrame) -> str:
    if dataframe.empty:
        return ""

    frame = dataframe.fillna("")
    columns = [str(column).strip() or f"column_{index + 1}" for index, column in enumerate(frame.columns)]
    lines = [" | ".join(columns)]

    for row in frame.itertuples(index=False, name=None):
        values = [str(value).replace("\n", " ").strip() for value in row]
        if any(values):
            lines.append(" | ".join(values))

    return "\n".join(lines)


def _export_google_file(service: Any, file_id: str, export_mime_type: str) -> bytes:
    request = service.files().export_media(fileId=file_id, mimeType=export_mime_type)
    buffer = io.BytesIO()
    downloader = MediaIoBaseDownload(buffer, request)

    done = False
    while not done:
        _, done = downloader.next_chunk()

    return buffer.getvalue()


def read_pdf(file_path: str | Path) -> str:
    """Extract text from a PDF using PyMuPDF."""
    try:
        import fitz
    except ImportError:
        raise ImportError("PDF support requires PyMuPDF (fitz).")

    path = _ensure_path(file_path)
    pages: list[str] = []

    with fitz.open(path) as document:
        for page in document:
            page_text = page.get_text("text")
            if page_text:
                pages.append(page_text)

    return _normalize_text("\n\n".join(pages))


def read_doc(file_path: str | Path) -> str:
    """Extract text from a DOC file using Microsoft Word automation on Windows."""
    path = _ensure_path(file_path)

    if win32com is not None:
        try:
            global word_app
            if word_app is None:
                word_app = win32com.client.Dispatch("Word.Application")
                word_app.Visible = False
            
            doc = word_app.Documents.Open(str(path))
            text = doc.Content.Text
            doc.Close()
            return _normalize_text(text)
        except Exception:
            pass

    return ""


def read_docx(file_path: str | Path) -> str:
    """Extract paragraphs and table content from a DOCX document."""
    path = _ensure_path(file_path)

    try:
        from docx import Document
        document = Document(path)
        blocks: list[str] = []

        for paragraph in document.paragraphs:
            if paragraph.text.strip():
                blocks.append(paragraph.text)

        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    blocks.append(" | ".join(cells))

        text = "\n".join(blocks)
        if text.strip():
            return _normalize_text(text)
    except Exception as exc:
        if docx2txt is None:
            raise ImportError("DOCX support requires python-docx or docx2txt.") from exc

    if docx2txt is None:
        raise ImportError("DOCX support requires python-docx or docx2txt.")

    text = docx2txt.process(str(path))
    return _normalize_text(text)


def read_txt(file_path: str | Path) -> str:
    """Extract text from a UTF-8 compatible plain text file."""
    return _normalize_text(_read_text_file(file_path))


def read_csv(file_path: str | Path) -> str:
    """Extract text from a CSV file using pandas."""
    path = _ensure_path(file_path)
    dataframe = pd.read_csv(path, dtype=str, keep_default_na=False, encoding="utf-8")
    return _normalize_text(_dataframe_to_text(dataframe))


def read_excel(file_path: str | Path) -> str:
    """Extract text from all sheets in an Excel workbook."""
    path = _ensure_path(file_path)
    engine = "xlrd" if path.suffix.lower() == ".xls" else "openpyxl"
    sheets = pd.read_excel(path, sheet_name=None, dtype=str, engine=engine)
    blocks: list[str] = []

    for sheet_name, dataframe in sheets.items():
        sheet_text = _dataframe_to_text(dataframe)
        if sheet_text:
            blocks.append(f"Sheet: {sheet_name}\n{sheet_text}")

    return _normalize_text("\n\n".join(blocks))


def read_pptx(file_path: str | Path) -> str:
    """Extract text from all text-bearing shapes in a PowerPoint presentation."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("PPTX support requires python-pptx.")

    path = _ensure_path(file_path)
    presentation = Presentation(path)
    blocks: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = [f"Slide {index}"]
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                slide_lines.append(text)
        blocks.append("\n".join(slide_lines))

    return _normalize_text("\n\n".join(blocks))


def read_html(file_path: str | Path) -> str:
    """Extract visible text from an HTML document."""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("HTML support requires beautifulsoup4.")

    soup = BeautifulSoup(_read_text_file(file_path), "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return _normalize_text(soup.get_text(separator="\n"))


def read_json(file_path: str | Path) -> str:
    """Extract text from a JSON file by pretty-printing it as UTF-8 text."""
    path = _ensure_path(file_path)
    with path.open("r", encoding="utf-8", errors="ignore") as handle:
        payload = json.load(handle)
    return _normalize_text(json.dumps(payload, ensure_ascii=False, indent=2))


def read_markdown(file_path: str | Path) -> str:
    """Extract readable text from a Markdown file via HTML conversion."""
    try:
        from bs4 import BeautifulSoup
        from markdown import markdown
    except ImportError:
        raise ImportError("Markdown support requires beautifulsoup4 and Markdown.")

    html = markdown(_read_text_file(file_path), output_format="html5")
    soup = BeautifulSoup(html, "html.parser")
    return _normalize_text(soup.get_text(separator="\n"))


def read_xml(file_path: str | Path) -> str:
    """Extract text content from an XML document."""
    path = _ensure_path(file_path)
    root = ElementTree.parse(path).getroot()
    parts = [text.strip() for text in root.itertext() if text and text.strip()]
    return _normalize_text("\n".join(parts))


def read_log(file_path: str | Path) -> str:
    """Extract text from a log file."""
    return _normalize_text(_read_text_file(file_path))


def read_rtf(file_path: str | Path) -> str:
    """Extract text from an RTF file."""
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        raise ImportError("RTF support requires striprtf.")

    return _normalize_text(rtf_to_text(_read_text_file(file_path)))


def read_odt(file_path: str | Path) -> str:
    """Extract text from an ODT document."""
    if load_odt is None or teletype is None or P is None or H is None:
        raise ImportError("ODT support requires odfpy.")

    path = _ensure_path(file_path)
    document = load_odt(str(path))
    elements = document.getElementsByType(H) + document.getElementsByType(P)
    parts = [teletype.extractText(element).strip() for element in elements]
    return _normalize_text("\n".join(part for part in parts if part))


def read_epub(file_path: str | Path) -> str:
    """Extract text from an EPUB file."""
    from bs4 import BeautifulSoup

    if epub is None:
        raise ImportError("EPUB support requires ebooklib.")

    path = _ensure_path(file_path)
    book = epub.read_epub(str(path))
    parts: list[str] = []

    for item in book.get_items():
        if item.get_type() == epub.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.get_content(), "html.parser")
            text = soup.get_text(separator="\n")
            if text.strip():
                parts.append(text)

    return _normalize_text("\n\n".join(parts))


def read_image_ocr(file_path: str | Path) -> str:
    """Extract text from an image using Tesseract OCR."""
    try:
        import pytesseract
        from PIL import Image, ImageOps
    except ImportError:
        raise ImportError("Image OCR support requires pytesseract and Pillow.")

    path = _ensure_path(file_path)
    with Image.open(path) as image:
        processed = ImageOps.grayscale(image)
        text = pytesseract.image_to_string(processed)
    return _normalize_text(text)


def read_google_doc(service: Any, file_id: str) -> str:
    """Export a Google Doc as plain text and return normalized content."""
    content = _export_google_file(service, file_id, "text/plain")
    return _normalize_text(content.decode("utf-8", errors="ignore"))


def read_google_sheet(service: Any, file_id: str) -> str:
    """Export a Google Sheet as CSV and return normalized content."""
    content = _export_google_file(service, file_id, "text/csv")
    dataframe = pd.read_csv(io.StringIO(content.decode("utf-8", errors="ignore")), dtype=str, keep_default_na=False)
    return _normalize_text(_dataframe_to_text(dataframe))


def read_google_slide(service: Any, file_id: str) -> str:
    """Export a Google Slides file as text and return normalized content."""
    content = _export_google_file(service, file_id, "text/plain")
    return _normalize_text(content.decode("utf-8", errors="ignore"))


def extract_text(file_path: str | Path) -> str:
    """Detect a local file type by extension and extract normalized text."""
    path = _ensure_path(file_path)
    reader = SUPPORTED_EXTENSIONS.get(path.suffix.lower())
    if reader is None:
        raise ValueError(f"Unsupported file type: {path.suffix or '<none>'}")
    return reader(path)


def extract_google_file(service: Any, file_id: str, mime_type: str) -> str:
    """Extract normalized text from a supported Google Drive native file."""
    mime_type = mime_type.strip().lower()

    if mime_type == GOOGLE_DOC_MIME:
        return read_google_doc(service, file_id)
    if mime_type == GOOGLE_SHEET_MIME:
        return read_google_sheet(service, file_id)
    if mime_type == GOOGLE_SLIDE_MIME:
        return read_google_slide(service, file_id)

    raise ValueError(f"Unsupported Google Drive mime type: {mime_type}")
